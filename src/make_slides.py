"""
Generates the final presentation PPTX for the housing analysis project.
Run from the project root: python src/make_slides.py
Output: outputs/housing_analysis.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # near-black
RED      = RGBColor(0xCC, 0x00, 0x00)   # TTU scarlet
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xE2, 0xE8, 0xF0)   # body text
MGRAY    = RGBColor(0x94, 0xA3, 0xB8)   # muted text
CARD     = RGBColor(0x1E, 0x28, 0x3A)   # card background
GREEN    = RGBColor(0x4A, 0xDE, 0x80)
ORANGE   = RGBColor(0xFF, 0xA5, 0x00)

# ── Canvas ────────────────────────────────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)
HH = Inches(0.9)          # header height

os.makedirs("outputs", exist_ok=True)
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════════
# Primitive helpers
# ═══════════════════════════════════════════════════════════════════════════════

def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return tb

def bullets(slide, items, x, y, w, h, size=20, color=LGRAY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(size); r.font.color.rgb = color
    return tb

def header(slide, title):
    rect(slide, 0, 0, W, HH, RED)
    txt(slide, title,
        Inches(0.4), Inches(0.12), Inches(12.5), HH,
        size=32, bold=True, color=WHITE)

def footer(slide, text="ISQS 3358  ·  Spring 2026  ·  Texas Tech University"):
    txt(slide, text,
        Inches(0.3), Inches(7.18), Inches(12.7), Inches(0.28),
        size=11, color=MGRAY, align=PP_ALIGN.CENTER)

def img(slide, path, x, y, w, h=None):
    if not os.path.exists(path): return
    if h is None: slide.shapes.add_picture(path, x, y, width=w)
    else:         slide.shapes.add_picture(path, x, y, width=w, height=h)

def label(slide, text, x, y):
    txt(slide, text.upper(), x, y, Inches(8), Inches(0.28),
        size=11, bold=True, color=RED)

# Content area below header: top = HH + gap, usable height ≈ 6.3"
CT = HH + Inches(0.1)   # content top
CB = Inches(7.1)         # content bottom (above footer)
CH = CB - CT             # ≈ 6.1"
CL = Inches(0.4)         # content left
CR = W - Inches(0.4)    # content right
CW = CR - CL             # ≈ 12.53"


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, W, Inches(0.3), RED)
rect(s, 0, Inches(7.2), W, Inches(0.3), RED)

txt(s, "How Does Economic Performance\nAffect the Housing Market?",
    Inches(0.8), Inches(0.8), Inches(11.7), Inches(2.6),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "A Regression Study on Real U.S. Home Prices  ·  1984 – 2025",
    Inches(0.8), Inches(3.55), Inches(11.7), Inches(0.55),
    size=20, color=LGRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(3.8), Inches(4.22), Inches(5.7), Inches(0.05), RED)

txt(s, "Samuel Melghem   ·   Omar Hernandez   ·   Ryan Rhodes\n"
       "Ethan Hennenhoefer   ·   Logan Crider",
    Inches(0.8), Inches(4.42), Inches(11.7), Inches(1.0),
    size=18, color=LGRAY, align=PP_ALIGN.CENTER)

txt(s, "ISQS 3358  ·  Spring 2026  ·  Texas Tech University",
    Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.45),
    size=15, color=MGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Roadmap  (4 + 4 layout, larger boxes)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s); header(s, "Presentation Roadmap"); footer(s)

items = [
    ("1", "Project Purpose",         "Why housing? Why it matters"),
    ("2", "Industry Background",      "The housing market at a glance"),
    ("3", "Datasets & Sources",       "7 FRED series, 1984 – 2025"),
    ("4", "Data Pipeline",            "Download → clean → merge → model"),
    ("5", "Exploratory Analysis",     "Trends, relationships, affordability"),
    ("6", "Regression Model",         "Features, methodology, results"),
    ("7", "Key Findings",             "What drives home prices"),
    ("8", "Limitations & Conclusion", "What the model can and can't do"),
]

BW = Inches(2.95); BH = Inches(1.45); GAP = Inches(0.22)
for i, (num, title, sub) in enumerate(items):
    col = i % 4; row = i // 4
    x = CL + col * (BW + GAP)
    y = CT + Inches(0.15) + row * (BH + Inches(0.35))
    rect(s, x, y, BW, BH, CARD)
    rect(s, x, y, Inches(0.3), BH, RED)
    txt(s, num,  x + Inches(0.05), y + Inches(0.22), Inches(0.28), Inches(0.7),
        size=26, bold=True, color=WHITE)
    txt(s, title, x + Inches(0.38), y + Inches(0.08), BW - Inches(0.45), Inches(0.5),
        size=14, bold=True, color=WHITE, wrap=True)
    txt(s, sub,   x + Inches(0.38), y + Inches(0.62), BW - Inches(0.45), Inches(0.75),
        size=12, color=MGRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — Project Purpose
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s); header(s, "Project Purpose"); footer(s)

txt(s, "Housing is the single largest financial decision most Americans make —\n"
       "yet few understand what actually drives prices.",
    CL, CT + Inches(0.05), CW, Inches(0.9),
    size=21, color=LGRAY)

bullets(s, [
    "Median home price rose from $79K in 1984 to $420K+ by 2025 — a 5× increase in nominal terms",
    "Inflation-adjusted, real prices roughly doubled — but income growth has lagged behind",
    "Mortgage rates swung from 14% (1984) to under 3% (2021), reshaping who can afford to buy",
    "The 2008 crisis showed that mispriced housing can take down the entire economy",
],
    CL, CT + Inches(1.05), CW, Inches(3.2), size=20)

rect(s, CL, CT + Inches(4.45), CW, Inches(1.1), CARD)
rect(s, CL, CT + Inches(4.45), Inches(0.07), Inches(1.1), RED)
txt(s, "Research Question",
    CL + Inches(0.2), CT + Inches(4.5), CW - Inches(0.3), Inches(0.32),
    size=13, bold=True, color=RED)
txt(s, "Do economic indicators predict inflation-adjusted home prices — "
       "and which factors have the strongest effect?",
    CL + Inches(0.2), CT + Inches(4.85), CW - Inches(0.3), Inches(0.65),
    size=19, bold=True, color=WHITE, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — Industry Background
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s); header(s, "Industry Background"); footer(s)

LW = Inches(5.9); RW = Inches(6.2); GAP2 = Inches(0.43)
RX = CL + LW + GAP2

# Left
label(s, "The U.S. Housing Market", CL, CT + Inches(0.05))
bullets(s, [
    "$40+ trillion in U.S. residential real estate value",
    "Prices set by supply, demand, and financing costs",
    "Highly sensitive to Federal Reserve rate policy",
    "Only ~1.5 million new homes built per year — supply is slow",
], CL, CT + Inches(0.38), LW, Inches(2.5), size=19)

label(s, "Why It's Hard to Predict", CL, CT + Inches(2.95))
bullets(s, [
    "Mortgage rates move weekly; income data updates annually",
    "Supply shocks (COVID, 2008) create structural breaks",
    "Regional markets often diverge sharply from national averages",
], CL, CT + Inches(3.28), LW, Inches(2.0), size=19)

# Right — dark card
rect(s, RX, CT, RW, CH - Inches(0.1), CARD)
rect(s, RX, CT, RW, Inches(0.06), RED)
label(s, "Our Approach", RX + Inches(0.25), CT + Inches(0.12))
bullets(s, [
    "Pull 7 economic series from FRED (Federal Reserve Economic Data)",
    "Standardize all frequencies to monthly, merge into one dataset",
    "Inflation-adjust home prices: RealPrice = HomePrice ÷ CPI × 100",
    "Fit a linear regression model to quantify each factor's effect",
    "Full dataset: January 1984 – December 2025",
], RX + Inches(0.25), CT + Inches(0.5), RW - Inches(0.4), Inches(5.0), size=19)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — Datasets
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Datasets  ·  All from FRED (Federal Reserve Economic Data)"); footer(s)

rows = [
    ("Home Price",   "MSPUS",          "Median U.S. home sale price",           "Quarterly"),
    ("CPI",          "CPIAUCSL",       "Consumer Price Index (inflation base)",  "Monthly"),
    ("Mortgage Rate","MORTGAGE30US",   "30-year fixed mortgage rate",           "Weekly → Monthly"),
    ("Unemployment", "UNRATE",         "U.S. unemployment rate",                "Monthly"),
    ("Med. Income",  "MEHOINUSA646N",  "Median household income",               "Annual → Monthly"),
    ("Population",   "POPTHM",         "U.S. total population",                 "Monthly"),
    ("Housing Starts","HOUST",         "New housing construction starts",        "Monthly"),
]
hdrs  = ["Variable", "FRED ID", "Description", "Frequency"]
col_x = [CL, CL+Inches(1.95), CL+Inches(4.1), CL+Inches(11.05)]
col_w = [Inches(1.85), Inches(2.05), Inches(6.85), Inches(1.85)]
RH    = Inches(0.62)
Y0    = CT + Inches(0.05)

rect(s, CL, Y0, CW, RH, RED)
for hdr, cx, cw in zip(hdrs, col_x, col_w):
    txt(s, hdr, cx+Inches(0.1), Y0+Inches(0.1), cw, RH,
        size=15, bold=True, color=WHITE)

for i, (var, fid, desc, freq) in enumerate(rows):
    y = Y0 + (i+1)*RH
    rect(s, CL, y, CW, RH, CARD if i%2==0 else BG)
    for val, cx, cw in zip([var, fid, desc, freq], col_x, col_w):
        txt(s, val, cx+Inches(0.1), y+Inches(0.1), cw, RH, size=15, color=LGRAY)

yn = Y0 + 8*RH + Inches(0.08)
rect(s, CL, yn, CW, Inches(0.62), CARD)
txt(s, "Target:  RealPrice = HomePrice ÷ CPI × 100   (inflation-adjusted price)",
    CL+Inches(0.2), yn+Inches(0.1), CW-Inches(0.3), Inches(0.45),
    size=16, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Data Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Data Pipeline  ·  One Command Runs Everything"); footer(s)

steps = [
    ("1  Download",  "fetch_fred.py\n7 FRED series\n→ data/raw/"),
    ("2  Resample",  "resampler.py\nWeekly → monthly mean\nAnnual → forward-fill"),
    ("3  Merge",     "merge_data.py\nOuter join on date\nFill any gaps"),
    ("4  Adjust",    "calculate_real_price.py\nRealPrice =\nHomePrice ÷ CPI × 100"),
    ("5  Model",     "model.py\nLinear regression\n1984 – 2025 full fit"),
]
PBW = Inches(2.35); PBH = Inches(2.2); PBY = CT + Inches(0.1)
for i, (title, body) in enumerate(steps):
    x = CL + i*(PBW + Inches(0.18))
    rect(s, x, PBY, PBW, PBH, CARD)
    rect(s, x, PBY, PBW, Inches(0.42), RED)
    txt(s, title, x+Inches(0.12), PBY+Inches(0.06), PBW-Inches(0.2), Inches(0.35),
        size=14, bold=True, color=WHITE)
    txt(s, body,  x+Inches(0.12), PBY+Inches(0.5), PBW-Inches(0.2), PBH-Inches(0.58),
        size=14, color=LGRAY, wrap=True)

label(s, "Key Cleaning Decisions", CL, CT + Inches(2.55))
bullets(s, [
    "Weekly mortgage rate → averaged to monthly mean",
    "Annual income → forward-filled (same value held month-to-month until next release)",
    "All 7 series merged on date with outer join; remaining gaps filled forward",
    "Result: 504 clean monthly observations, January 1984 – December 2025",
], CL, CT + Inches(2.88), CW, Inches(2.8), size=19)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Individual Trends  (2 × 3 chart grid)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Exploratory Analysis  ·  Individual Trends (1984 – 2025)"); footer(s)

charts = [
    "outputs/charts/real_price.png",
    "outputs/charts/mortgage_rate.png",
    "outputs/charts/unemployment_rate.png",
    "outputs/charts/median_income.png",
    "outputs/charts/housing_starts.png",
    "outputs/charts/population.png",
]
IW = Inches(4.1); IH = Inches(2.7); IGAP = Inches(0.16)
for i, path in enumerate(charts):
    col = i % 3; row = i // 3
    x = CL + col*(IW + IGAP)
    y = CT + Inches(0.05) + row*(IH + Inches(0.2))
    img(s, path, x, y, IW, IH)

txt(s, "Real prices peaked in 2006, crashed in 2008, recovered and surged through 2025.  "
       "Mortgage rates fell from 14% → 3% before rebounding in 2022.",
    CL, Inches(7.0), CW, Inches(0.3),
    size=13, color=MGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Relationships  (2 key charts + callout)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Exploratory Analysis  ·  Key Relationships"); footer(s)

LCW = Inches(6.5); LCH = Inches(2.9)
img(s, "outputs/charts/price_vs_mortgage.png",    CL, CT + Inches(0.05), LCW, LCH)
img(s, "outputs/charts/price_vs_income.png",       CL, CT + Inches(3.1),  LCW, LCH)

RX2 = CL + LCW + Inches(0.35)
RW2 = W - RX2 - Inches(0.3)
rect(s, RX2, CT, RW2, CH - Inches(0.1), CARD)
rect(s, RX2, CT, RW2, Inches(0.06), RED)
label(s, "What We See", RX2+Inches(0.22), CT+Inches(0.12))
bullets(s, [
    "Mortgage rate (top): strong negative relationship — "
      "as rates fell from 14% → 3%, prices surged. "
      "Higher rates price buyers out of the market.",
    "Income vs Price (bottom, indexed to 100 in 1984): "
      "home prices grew nearly twice as fast as incomes — "
      "the gap widened dramatically after 2000 and again after 2020.",
],  RX2+Inches(0.22), CT+Inches(0.5), RW2-Inches(0.35), Inches(5.3), size=18)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Affordability
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Affordability  ·  Are Homes Getting Harder to Buy?"); footer(s)

AIW = Inches(6.15); AIH = Inches(3.1)
img(s, "outputs/charts/nominal_vs_real.png",    CL,             CT + Inches(0.05), AIW, AIH)
img(s, "outputs/charts/price_to_income_ratio.png", CL+AIW+Inches(0.25), CT + Inches(0.05), AIW, AIH)

rect(s, CL,                  CT+AIH+Inches(0.2), AIW,             Inches(2.55), CARD)
rect(s, CL+AIW+Inches(0.25), CT+AIH+Inches(0.2), AIW,             Inches(2.55), CARD)

txt(s, "Nominal vs. Real Price",
    CL+Inches(0.18), CT+AIH+Inches(0.28), AIW-Inches(0.3), Inches(0.35),
    size=16, bold=True, color=RED)
bullets(s, [
    "Nominal price rose 5× from 1984 to 2025",
    "Inflation-adjusted real price rose only ~2×",
    "Inflation explains more than half the apparent increase",
], CL+Inches(0.18), CT+AIH+Inches(0.68), AIW-Inches(0.3), Inches(1.9), size=17)

txt(s, "Price-to-Income Ratio",
    CL+AIW+Inches(0.43), CT+AIH+Inches(0.28), AIW-Inches(0.3), Inches(0.35),
    size=16, bold=True, color=RED)
bullets(s, [
    "1984: a home cost ~3.5 years of median income",
    "2025: that figure is now 7+ years — doubled",
    "Affordability has significantly deteriorated",
], CL+AIW+Inches(0.43), CT+AIH+Inches(0.68), AIW-Inches(0.3), Inches(1.9), size=17)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Model Setup
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Regression Model  ·  Setup & Methodology"); footer(s)

LW2 = Inches(6.0); RX3 = CL + LW2 + Inches(0.5); RW3 = CW - LW2 - Inches(0.5)

label(s, "Model Type", CL, CT+Inches(0.05))
txt(s, "Ordinary Least Squares Linear Regression",
    CL, CT+Inches(0.38), LW2, Inches(0.45), size=20, color=LGRAY)

label(s, "Features  (X — independent variables)", CL, CT+Inches(0.95))
bullets(s, [
    "Unemployment Rate",
    "Mortgage Rate (30-year fixed)",
    "Median Household Income",
    "U.S. Population",
    "Housing Starts",
], CL, CT+Inches(1.28), LW2, Inches(2.6), size=20)

label(s, "Target  (y — dependent variable)", CL, CT+Inches(3.98))
txt(s, "RealPrice  —  inflation-adjusted median home sale price",
    CL, CT+Inches(4.3), LW2, Inches(0.45), size=20, color=LGRAY)

label(s, "Why Linear Regression?", CL, CT+Inches(4.9))
bullets(s, [
    "Interpretable coefficients — we can say exactly how much each variable moves price",
    "Strong baseline for economic data with expected linear relationships",
], CL, CT+Inches(5.22), LW2, Inches(1.1), size=18)

# Right — dark card
rect(s, RX3, CT, RW3, CH-Inches(0.1), CARD)
rect(s, RX3, CT, RW3, Inches(0.06), RED)
label(s, "Data Coverage", RX3+Inches(0.22), CT+Inches(0.12))
txt(s, "January 1984 – December 2025",
    RX3+Inches(0.22), CT+Inches(0.45), RW3-Inches(0.35), Inches(0.5),
    size=22, bold=True, color=WHITE)
txt(s, "504 monthly observations\nAll data used for model fitting\n(no train/test split —\nfull dataset regression)",
    RX3+Inches(0.22), CT+Inches(1.05), RW3-Inches(0.35), Inches(2.4),
    size=18, color=LGRAY, wrap=True)

rect(s, RX3+Inches(0.22), CT+Inches(3.6), RW3-Inches(0.44), Inches(0.06), RED)
txt(s, "R²  =  0.8775",
    RX3+Inches(0.22), CT+Inches(3.8), RW3-Inches(0.35), Inches(0.7),
    size=32, bold=True, color=WHITE)
txt(s, "Model explains 88% of variation\nin real home prices",
    RX3+Inches(0.22), CT+Inches(4.55), RW3-Inches(0.35), Inches(0.8),
    size=17, color=LGRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 11 — Results
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Regression Results  ·  Coefficients & Model Performance"); footer(s)

# Top stat boxes
stats = [
    ("R²",          "0.8775",       "Explains 88% of variation\nin real home prices"),
    ("√MSE",        "≈ $6,546",     "Average prediction error\nin today's dollars"),
    ("Observations","504 months",   "January 1984\nthrough December 2025"),
    ("Features",    "5 variables",  "All FRED series\nafter inflation adjustment"),
]
SW = Inches(2.95); SH = Inches(1.5); SY = CT + Inches(0.05)
for i, (lbl, val, sub) in enumerate(stats):
    x = CL + i*(SW + Inches(0.22))
    rect(s, x, SY, SW, SH, CARD)
    rect(s, x, SY, SW, Inches(0.07), RED)
    txt(s, lbl, x+Inches(0.15), SY+Inches(0.1), SW-Inches(0.25), Inches(0.3),
        size=13, color=MGRAY)
    txt(s, val, x+Inches(0.15), SY+Inches(0.4), SW-Inches(0.25), Inches(0.55),
        size=26, bold=True, color=WHITE)
    txt(s, sub, x+Inches(0.15), SY+Inches(0.98), SW-Inches(0.25), Inches(0.45),
        size=12, color=MGRAY, wrap=True)

# Coefficient table
coefs = [
    ("Mortgage Rate",          "−847",    "↓ Negative", "Higher rates → fewer buyers → lower prices"),
    ("Unemployment Rate",      "−1,338",  "↓ Negative", "Higher unemployment → less buying power → lower prices"),
    ("Median Household Income","+0.55",   "↑ Positive", "More income → more purchasing power → higher prices"),
    ("Population",             "+0.20",   "↑ Positive", "More people → more demand → higher prices"),
    ("Housing Starts",         "+7.53",   "↑ Positive", "More construction during growth periods"),
]
hdrs2  = ["Feature", "Coefficient", "Direction", "Interpretation"]
cx2    = [CL, CL+Inches(3.05), CL+Inches(4.55), CL+Inches(5.85)]
cw2    = [Inches(2.95), Inches(1.4), Inches(1.25), Inches(7.1)]
TRH    = Inches(0.56)
TY0    = SY + SH + Inches(0.18)

rect(s, CL, TY0, CW, TRH, RED)
for h, cx, cw in zip(hdrs2, cx2, cw2):
    txt(s, h, cx+Inches(0.1), TY0+Inches(0.1), cw, TRH,
        size=14, bold=True, color=WHITE)

for i, (feat, coef, dirn, interp) in enumerate(coefs):
    y = TY0 + (i+1)*TRH
    rect(s, CL, y, CW, TRH, CARD if i%2==0 else BG)
    dc = GREEN if "Positive" in dirn else ORANGE
    for val, cx, cw, c in zip([feat, coef, dirn, interp],
                               cx2, cw2,
                               [LGRAY, WHITE, dc, LGRAY]):
        txt(s, val, cx+Inches(0.1), y+Inches(0.08), cw, TRH, size=14, color=c)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 12 — Actual vs Predicted
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "How Well Did the Model Predict?  ·  Actual vs. Predicted"); footer(s)

img(s, "outputs/charts/actual_vs_predicted.png",
    CL, CT+Inches(0.05), Inches(8.1), Inches(5.9))

RX4 = CL + Inches(8.4); RW4 = W - RX4 - Inches(0.3)
rect(s, RX4, CT, RW4, CH-Inches(0.1), CARD)
rect(s, RX4, CT, RW4, Inches(0.06), RED)
label(s, "Reading the Chart", RX4+Inches(0.22), CT+Inches(0.12))
bullets(s, [
    "Blue = actual inflation-adjusted price",
    "Dashed = model's fitted values",
], RX4+Inches(0.22), CT+Inches(0.45), RW4-Inches(0.35), Inches(1.3), size=17)

label(s, "What This Tells Us", RX4+Inches(0.22), CT+Inches(1.9))
bullets(s, [
    "Model tracks the 40-year trend extremely well (R² = 0.88)",
    "Captures the 2008 crash and the long recovery",
    "The 2020–2025 COVID surge is visible — the model follows it but economic fundamentals were stress-tested",
    "Short-term shocks move prices above or below the economic baseline",
], RX4+Inches(0.22), CT+Inches(2.25), RW4-Inches(0.35), Inches(3.6), size=16)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 13 — Key Findings  (3 × 2 cards, bigger)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s); header(s, "Key Findings"); footer(s)

cards = [
    ("Income Drives Prices Up",
     "Median household income had the strongest positive effect. "
     "As earnings rise, buyers qualify for larger mortgages — prices follow."),
    ("Mortgage Rates Suppress Prices",
     "Every 1% rise in the 30-year rate corresponds to ~$847 lower real price. "
     "The 40-year rate decline quietly subsidized price growth."),
    ("Unemployment Drives Prices Down",
     "Higher unemployment reduces buying power and forces distressed sellers — "
     "seen clearly in the 2008–2009 price crash."),
    ("Affordability Has Deteriorated",
     "Price-to-income ratio doubled from 3.5× in 1984 to 7×+ in 2025. "
     "Even adjusting for inflation, homes are much harder to afford."),
    ("Inflation Masks Real Growth",
     "Nominal prices rose 5×, but inflation-adjusted prices only ~2×. "
     "Inflation explains most of the apparent price explosion."),
    ("Model Explains 88% of Variation",
     "Five economic variables — income, rates, unemployment, population, "
     "and housing starts — capture the vast majority of price movement."),
]
FCW = Inches(4.0); FCH = Inches(2.5); FGAP = Inches(0.25)
for i, (title, body) in enumerate(cards):
    col = i % 3; row = i // 3
    x = CL + col*(FCW + FGAP)
    y = CT + Inches(0.08) + row*(FCH + Inches(0.18))
    rect(s, x, y, FCW, FCH, CARD)
    rect(s, x, y, Inches(0.07), FCH, RED)
    txt(s, title, x+Inches(0.18), y+Inches(0.1), FCW-Inches(0.25), Inches(0.48),
        size=16, bold=True, color=WHITE, wrap=True)
    txt(s, body,  x+Inches(0.18), y+Inches(0.62), FCW-Inches(0.25), FCH-Inches(0.72),
        size=14, color=LGRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 14 — Limitations
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Limitations  ·  What the Model Can't Fully Capture"); footer(s)

lims = [
    ("National Averages Only",
     "All FRED series are U.S.-wide. Texas, California, and rural Midwest "
     "behave very differently — a national model misses all regional variation."),
    ("Multicollinearity",
     "Income, population, and time all trend together. This can inflate or deflate "
     "individual coefficients — interpret directions, not exact magnitudes."),
    ("Linear Assumption",
     "The model assumes relationships are constant over 40 years. In reality, "
     "a 1% mortgage rate change at 14% differs from one at 3%."),
    ("Supply-Side Gaps",
     "Zoning laws, construction costs, and land scarcity matter. "
     "Housing starts is a rough proxy, not a full picture of supply."),
]
LRH = Inches(1.38)
for i, (title, body) in enumerate(lims):
    y = CT + Inches(0.05) + i*(LRH + Inches(0.12))
    rect(s, CL, y, CW, LRH, CARD)
    rect(s, CL, y, Inches(0.07), LRH, RED)
    txt(s, title, CL+Inches(0.2), y+Inches(0.1), Inches(2.6), Inches(0.45),
        size=16, bold=True, color=WHITE)
    txt(s, body,  CL+Inches(2.95), y+Inches(0.12), CW-Inches(3.05), LRH-Inches(0.25),
        size=17, color=LGRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 15 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s); header(s, "Conclusion  ·  What We Learned"); footer(s)

txt(s, "Economic performance drives the housing market — but not always in the ways people assume.",
    CL, CT+Inches(0.05), CW, Inches(0.6), size=21, color=LGRAY)

LCL = Inches(6.0); RCL = CW - LCL - Inches(0.35)
RXC = CL + LCL + Inches(0.35)

rect(s, CL,  CT+Inches(0.8), LCL, CH-Inches(0.95), CARD)
rect(s, CL,  CT+Inches(0.8), LCL, Inches(0.06), RED)
txt(s, "What We Found",
    CL+Inches(0.2), CT+Inches(0.9), LCL-Inches(0.3), Inches(0.4),
    size=16, bold=True, color=RED)
bullets(s, [
    "Income and mortgage rates are the two strongest forces on home prices",
    "Homes now cost 7× median annual income — up from 3.5× in 1984",
    "Inflation explains most of the 5× nominal price increase; real growth was ~2×",
    "Five economic variables explain 88% of 40 years of price variation",
    "Short-term shocks (2008 crisis, COVID) add noise above the economic baseline",
], CL+Inches(0.2), CT+Inches(1.35), LCL-Inches(0.3), Inches(4.3), size=18)

rect(s, RXC, CT+Inches(0.8), RCL, CH-Inches(0.95), CARD)
rect(s, RXC, CT+Inches(0.8), RCL, Inches(0.06), RED)
txt(s, "The Big Picture",
    RXC+Inches(0.2), CT+Inches(0.9), RCL-Inches(0.3), Inches(0.4),
    size=16, bold=True, color=RED)
bullets(s, [
    "The 40-year decline in mortgage rates (14% → 3%) quietly subsidized price growth — "
      "cheap financing made expensive homes feel affordable",
    "When rates reversed sharply in 2022, affordability collapsed almost overnight",
    "Data-driven analysis of fundamentals is essential — nominal prices alone are misleading",
], RXC+Inches(0.2), CT+Inches(1.35), RCL-Inches(0.3), Inches(3.0), size=18)

rect(s, CL, CT+Inches(5.25), CW, Inches(0.65), CARD)
txt(s, "github.com/EthanH2004/Housing_Analysis  "
       "·  Full code & data, one command reproduces the entire analysis",
    CL+Inches(0.2), CT+Inches(5.35), CW-Inches(0.3), Inches(0.5),
    size=15, color=MGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 16 — Q&A
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, W, Inches(0.3), RED)
rect(s, 0, Inches(7.2), W, Inches(0.3), RED)

txt(s, "Thank You",
    Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.4),
    size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Questions?",
    Inches(0.5), Inches(2.85), Inches(12.3), Inches(0.8),
    size=32, color=LGRAY, align=PP_ALIGN.CENTER)
rect(s, Inches(3.5), Inches(3.8), Inches(6.3), Inches(0.06), RED)
txt(s, "Samuel Melghem   ·   Omar Hernandez   ·   Ryan Rhodes\n"
       "Ethan Hennenhoefer   ·   Logan Crider",
    Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.0),
    size=19, color=LGRAY, align=PP_ALIGN.CENTER)
txt(s, "ISQS 3358  ·  Spring 2026  ·  Texas Tech University",
    Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.45),
    size=15, color=MGRAY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "outputs/housing_analysis.pptx"
prs.save(out)
print(f"\nSaved → {out}")
print(f"Total slides: {len(prs.slides)}")
